# pptxを扱うためのモジュールをインポートする
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

#-----------------------------------------------
# 関数：テキストボックスを描画する
#-----------------------------------------------
def add_textbox(
    slide,
    text,
    left=0, top=0,
    width=0, height=0,
    font_size=12,
    font_name="メイリオ",
    align="left",
    valign="top",
    font_color=(0, 0, 0),
    bg_color=None,
    border_color=None,
    border_width=0.75,
    border_style="solid",
):

    textbox = slide.shapes.add_textbox(
        Cm(left), Cm(top),
        Cm(width), Cm(height)
    )

    tf = textbox.text_frame

    # 安定化設定（重要）
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE

    # 余白ゼロ（ズレ対策）
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_left = 0
    tf.margin_right = 0

    # 縦位置
    if valign == "middle":
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    elif valign == "top":
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    else:
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM

    # テキスト
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text

    # 行間（固定設定）
    p.line_spacing=1.3

    # 横位置
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "left":
        p.alignment = PP_ALIGN.LEFT
    else:
        p.alignment = PP_ALIGN.RIGHT

    # フォント
    run.font.size = Pt(font_size)
    run.font.name = font_name
    run.font.color.rgb = RGBColor(*font_color)

    # 背景
    fill = textbox.fill
    if bg_color is None:
        fill.background()  # 透明
    else:
        fill.solid()
        fill.fore_color.rgb = RGBColor(*bg_color)

    # 枠線
    line = textbox.line

    if border_color is None:
        line.fill.background()  # 枠線なし
    else:

        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(*border_color)

        # 太さ
        line.width = Pt(border_width)

        # 種類
        if border_style == "dash":
            line.dash_style = MSO_LINE_DASH_STYLE.DASH
        elif border_style == "dot":
            line.dash_style = MSO_LINE_DASH_STYLE.DOT
        else:
            line.dash_style = MSO_LINE_DASH_STYLE.SOLID

    # 影を消す
    textbox.shadow.inherit = False
    textbox.shadow.visible = False

    return textbox

#-----------------------------------------------
# 関数：線を引く
#-----------------------------------------------
def add_line(
    slide,
    x1, y1, x2, y2,
    color=(0, 0, 0),
    width=0.75,
    style="solid",  # solid / dash / dot
):
    line_shape = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Cm(x1), Cm(y1),
        Cm(x2), Cm(y2)
    )

    # 影を消す
    line_shape.shadow.inherit = False
    line_shape.shadow.visible = False

    line = line_shape.line

    # 色
    line.color.rgb = RGBColor(*color)

    # 太さ
    line.width = Pt(width)

    # 線種
    if style == "dash":
        line.dash_style = MSO_LINE_DASH_STYLE.DASH
    elif style == "dot":
        line.dash_style = MSO_LINE_DASH_STYLE.DOT
    else:
        line.dash_style = MSO_LINE_DASH_STYLE.SOLID

    return line_shape

#-----------------------------------------------
# メイン処理
#-----------------------------------------------
prs = Presentation()

# A4横
prs.slide_width = Cm(29.7)
prs.slide_height = Cm(21.0)

# --- 1ページ目 空白スライド ---
slide = prs.slides.add_slide(prs.slide_layouts[6])

# --- 見出し ---
add_textbox(
    slide,
    "1-1. タイトル",
    left=1, top=0.5,
    height=1.5, width=27.7,
    font_size=24,
    font_color=(50, 50, 50),
)

# --- 見出しの下の線 ---
add_line(
    slide,
    x1=0.5, y1=2.2, x2=29.2, y2=2.2,
    color=(220, 220, 220),
    width=2,
)

# --- 本文 ---
add_textbox(
    slide,
    "これは本文です。少し長い文章も自動で折り返されます。これは本文です。少し長い文章も自動で折り返されます。これは本文です。少し長い文章も自動で折り返されます。",
    left=2.5, top=5, height=1.75, width=24.7,
    font_size=16,
    font_color=(50, 50, 50),
    border_color=(220, 220, 220),
)

# --- 2ページ目 空白スライド ---
slide = prs.slides.add_slide(prs.slide_layouts[6])

# --- 見出し ---
add_textbox(
    slide,
    "1-2. 2ページ目のタイトル",
    left=1, top=0.5,
    height=1.5, width=27.7,
    font_size=24,
    font_color=(50, 50, 50),
)

# 保存
prs.save("pptx_sample.pptx")
print("pptx_sample.pptx を作成しました.")
